#!/usr/bin/env python3
"""Generate UHG Risk Intelligence panel demo PowerPoint."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUTPUT = ROOT / "UHG_Risk_Intelligence_Panel_Demo.pptx"

NAVY = RGBColor(0x00, 0x33, 0x66)
TEAL = RGBColor(0x00, 0x7A, 0x87)
GRAY = RGBColor(0x55, 0x55, 0x55)


def _set_title(shape, text: str, size: int = 32) -> None:
    tf = shape.text_frame if hasattr(shape, "text_frame") else shape
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = NAVY


def _add_bullets(text_frame, items: list[str], size: int = 18, level0: bool = True) -> None:
    text_frame.clear()
    for i, item in enumerate(items):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = item
        p.level = 0 if level0 else 0
        p.font.size = Pt(size)
        p.font.color.rgb = GRAY
        p.space_after = Pt(6)


def _add_slide_title_only(prs: Presentation, title: str, subtitle: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(Inches(0.6), Inches(2.8), Inches(12), Inches(1.2))
    tf = box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(11.5), Inches(1.5))
        stf = sub.text_frame
        stf.text = subtitle
        sp = stf.paragraphs[0]
        sp.font.size = Pt(20)
        sp.font.color.rgb = TEAL
        sp.alignment = PP_ALIGN.CENTER


def _add_content_slide(prs: Presentation, title: str, bullets: list[str], note: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # title and content
    _set_title(slide.shapes.title, title, size=28)
    body = slide.placeholders[1].text_frame
    _add_bullets(body, bullets, size=17)
    if note:
        p = body.add_paragraph()
        p.text = note
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = TEAL


def _add_two_column_slide(
    prs: Presentation,
    title: str,
    left_title: str,
    left_items: list[str],
    right_title: str,
    right_items: list[str],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.8))
    _set_title(title_box, title, 28)

    left = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5))
    ltf = left.text_frame
    ltf.text = left_title
    ltf.paragraphs[0].font.bold = True
    ltf.paragraphs[0].font.size = Pt(18)
    ltf.paragraphs[0].font.color.rgb = TEAL
    for item in left_items:
        p = ltf.add_paragraph()
        p.text = item
        p.font.size = Pt(15)
        p.font.color.rgb = GRAY
        p.level = 0

    right = slide.shapes.add_textbox(Inches(6.6), Inches(1.2), Inches(5.8), Inches(5.5))
    rtf = right.text_frame
    rtf.text = right_title
    rtf.paragraphs[0].font.bold = True
    rtf.paragraphs[0].font.size = Pt(18)
    rtf.paragraphs[0].font.color.rgb = TEAL
    for item in right_items:
        p = rtf.add_paragraph()
        p.text = item
        p.font.size = Pt(15)
        p.font.color.rgb = GRAY


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 Title
    _add_slide_title_only(
        prs,
        "UHG Geopolitical & Regulatory\nRisk Intelligence — Panel Demo",
        "Proof of Concept  |  Governed pipeline, not a chatbot",
    )

    # 2 Problem
    _add_content_slide(
        prs,
        "The Problem",
        [
            "UHG operates across Optum Health, Optum Insight, Optum Rx, and related segments",
            "Shared disruptors: Medicare/regulatory changes, cyber/data risk, political volatility",
            "Analysts and executives need a consistent, auditable answer to:",
            "   “How exposed is segment X to disruptor Y — and what pathway should we consider?”",
            "Today: one-off spreadsheets or unconstrained LLM chats — neither is governed or reproducible",
        ],
    )

    # 3 Thesis
    _add_content_slide(
        prs,
        "What This POC Shows",
        [
            "A designed-but-unvalidated framework adapted to UnitedHealth Group",
            "Deterministic quant core produces scores; LLMs do narrow, bounded tasks only",
            "Fail-closed guardrails: refuse when out of scope, unauthorized, or low confidence",
            "Human-in-the-loop for high-judgment overlay extractions (always staged first)",
            "One-line thesis: governed risk intelligence — numbers from code, not from chat",
        ],
        note="Not production-ready. No proprietary UHG data used.",
    )

    # 4 Architecture
    _add_content_slide(
        prs,
        "Pipeline Architecture (LangGraph)",
        [
            "1. Natural language question → Query parser (LLM)",
            "2. Guardrails: UHG scope, retired segments → Access control (role × segment)",
            "3. Quant core (deterministic): driver baseline — risk, opportunity, coverage, confidence",
            "4. Scenario engine: disruptor-adjusted scores + pathway recommendation",
            "5. Narrative generator (LLM): analyst or executive audience",
            "",
            "Separate path: Overlay extraction (LLM + search/MCP) → STAGED → analyst confirms once",
            "Optional: MCP tools for supplemental context — never affects scenario scores",
        ],
    )

    # 5 What you're proving
    _add_content_slide(
        prs,
        "What the Panel Should Take Away",
        [
            "Separation of concerns — LLMs parse & narrate; scores come from deterministic code",
            "Governance & trust — guardrails block bad input, shaky recommendations, unauthorized views",
            "Production-minded design — LangGraph, Phoenix observability, 37 unit tests, role-based access",
            "Honest POC scope — real public anchor events + synthetic indicators; validation roadmap explicit",
        ],
    )

    # 6 Demo 1-2
    _add_two_column_slide(
        prs,
        "Live Demo — Core Query Path",
        "Happy path",
        [
            'python demo.py query "How exposed is Optum Health to a Medicare Advantage rate cut?"',
            "Shows: segment, driver, disruptor",
            "Deterministic risk / opportunity scores",
            "Recommended pathway + analyst narrative",
            "Numbers are reproducible; LLM only structures & summarizes",
        ],
        "Executive audience",
        [
            'python demo.py query "How exposed is Optum Insight to a cyber disruptor?" --audience executive',
            "Same pipeline, one-sentence executive summary",
            "Same scenario object — different presentation layer",
        ],
    )

    # 7 Guardrails demo
    _add_content_slide(
        prs,
        "Live Demo — Guardrails (Fail Closed)",
        [
            'Out of scope:  query "Should I buy NVIDIA stock?"  →  out_of_scope',
            'Retired segment:  "Optum International"  →  unknown_segment (no silent substitution)',
            'Access control:  --role guest + Optum Insight  →  access_denied (guest → Optum Rx only)',
            'Confidence floor:  --sparse-data Capital  →  no pathway; narrative shows caveat',
            "System refuses rather than force-fitting or over-confident recommendations",
        ],
    )

    # 8 Overlay
    _add_content_slide(
        prs,
        "Live Demo — Overlay Extraction (Human in the Loop)",
        [
            "python demo.py overlay --segment Optum_Insight --driver Data_Digital",
            "Optional: --live-search (Tavily) or --mcp-search",
            "Flow: text in → LLM rates severity / immediacy / persistence",
            "Status is always STAGED — never live until analyst confirms",
            "Second confirmation attempt is rejected (invalid_transition guardrail)",
            "Connects to real anchor: 2024 Change Healthcare cyberattack (Optum Insight)",
        ],
    )

    # 9 MCP
    _add_content_slide(
        prs,
        "Live Demo — MCP Supplemental Tools",
        [
            "python demo.py tool list",
            "python demo.py tool call local.get_backtest_anchors",
            "python demo.py tool call yfmcp.yfinance_get_ticker_info --symbol UNH",
            "python demo.py query \"...\" --mcp-context",
            "Formatted analyst / executive summaries — not raw JSON dumps",
            "Explicitly non-authoritative: does NOT change scenario scores",
            "Phoenix traces: mcp.call_tool spans when OBSERVABILITY_BACKEND=phoenix",
        ],
    )

    # 10 Guardrails table
    _add_content_slide(
        prs,
        "Guardrails Cheat Sheet (Q&A)",
        [
            "out_of_scope — no UHG/Optum markers → domain boundary",
            "unknown_segment — retired/divested names → no nearest-match substitution",
            "access_denied — role cannot view segment → authorization separate from trust",
            "confidence floor — thin data → no pathway recommendation",
            "overlay staging — always staged until human confirm → subjective LLM judgment gated",
            "invalid_transition — second overlay confirm blocked → state machine integrity",
        ],
    )

    # 11 Data honesty
    _add_two_column_slide(
        prs,
        "Data Provenance — Be Explicit with the Panel",
        "Real (public)",
        [
            "Segment revenue — UHG FY2025 SEC / earnings",
            "Backtest anchors — MA rate impact (Optum Health)",
            "Change Healthcare cyberattack (Optum Insight)",
        ],
        "Synthetic (illustrative)",
        [
            "Indicator registry & driver loadings",
            "Expert-elicited weights",
            "Sample overlay event texts (paraphrases)",
            "No proprietary or internal UHG data required",
        ],
    )

    # 12 10-min flow
    _add_content_slide(
        prs,
        "Suggested 10-Minute Panel Flow",
        [
            "1. Problem + thesis (90 sec)",
            "2. Happy-path query (2 min)",
            "3. One guardrail — out of scope or access denied (1 min)",
            "4. Confidence floor (1 min)",
            "5. Overlay + human confirm (2 min)",
            "6. MCP supplemental context (1 min)",
            "7. Phoenix observability screenshot (30 sec)",
            "8. Limitations + next steps: real feeds, eval datasets, overlay calibration (1 min)",
        ],
    )

    # 13 Next steps
    _add_content_slide(
        prs,
        "Validation Roadmap (Next Steps)",
        [
            "Query-parser evals: labeled NL questions → expected ScenarioQuery",
            "Overlay calibration: double-rated ground truth (two raters per event)",
            "Core backtest: would Capital/Data_Digital baselines have flagged anchor events early?",
            "Wire real CMS / regulatory / cyber indicator feeds (replace synthetic registry)",
            "LangSmith or Phoenix eval datasets for regression on guardrails + parser accuracy",
        ],
    )

    # 14 Closing
    _add_slide_title_only(
        prs,
        "Closing",
        "Deterministic scoring at the center · LLMs in narrow roles · Guardrails that fail closed\n"
        "Human confirmation where judgment matters · Clear path from POC to validation",
    )

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    path = build()
    print(f"Created: {path}")
